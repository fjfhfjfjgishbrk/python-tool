import cv2
import os
import numpy as np
import math
import pptx
from pptx.util import Inches, Pt
import time

GREYSCALE_THRESHOLD = 0
WIDTH_LIMIT = 160
FPS = 20
FONT_SIZE = 6
INVERSE = True


def rescaleImg(imgToScale):
    imgHeight, imgWidth = imgToScale.shape[:2]
    scaleFactor = WIDTH_LIMIT / imgWidth
    if scaleFactor > 1:
        scaleFactor = 1
    return cv2.resize(imgToScale, (math.floor(imgWidth * scaleFactor), math.floor(imgHeight * scaleFactor)))


def getFrame(sec, c):
    vidcap.set(cv2.CAP_PROP_POS_MSEC, round(sec*1000))
    hasFrames, image = vidcap.read()
    if hasFrames:
        cv2.imwrite("frames/frame%d.jpg" % c, rescaleImg(image))  # save frame as JPEG file
    return hasFrames


filelist = [f for f in os.listdir("frames")]
for f in filelist:
    os.remove(os.path.join("frames", f))

vidcap = cv2.VideoCapture('video.mp4')
success, image = vidcap.read()
count = 0
sec = 0
secAdd = 1 / FPS
while success:
    sec += secAdd
    success = getFrame(sec, count)
    if count % 100 == 0:
        print('Read', count, "frames,", round(sec, 2), "sec")
    count += 1

count -= 1
print("Total frames:", count)
prs = pptx.Presentation()
greyMean = 0

for frame in range(count):
    if frame == 0:
        startTime = time.time()

    img = cv2.imread("frames/frame%d.jpg" % frame, cv2.IMREAD_GRAYSCALE)
    if INVERSE:
        for i in img:
            for j in i:
                j = 256 - j
    addVal = [1, 2, 4, 64, 8, 16, 32, 128]
    ascii_out = np.zeros((len(img) // 4 + 1, len(img[0]) // 2 + 1))
    x = 0
    imgMean = np.mean(img)
    greyMean = imgMean if greyMean == 0 or abs(greyMean - imgMean) > 15 else greyMean
    GREYSCALE_THRESHOLD = min(greyMean - 10, 230)
    for x_row in img:
        y = 0
        for scale in x_row:
            if int(scale) > GREYSCALE_THRESHOLD:
                ascii_out[x // 4][y // 2] += addVal[(x % 4) + 4 * (y % 2)]
            elif y % 2 == 0:
                if int(scale) + int(img[x][y + 1]) > GREYSCALE_THRESHOLD > int(img[x][y + 1]):
                    ascii_out[x // 4][y // 2] += addVal[x % 2 * 4 + x % 4]
            elif y % 2 == 0 and x % 2 == 0:
                if int(scale + img[x + 1][y] + img[x][y + 1] + img[x + 1][y + 1]) > GREYSCALE_THRESHOLD:
                    if not(int(scale) + int(img[x][y + 1]) > GREYSCALE_THRESHOLD or int(img[x + 1][y]) + int(img[x + 1][y + 1]) > GREYSCALE_THRESHOLD):
                        ascii_out[x // 4][y // 2] += addVal[5 + x % 4]
            y += 1
        x += 1

    # print(ascii_out)

    text_out = ""
    for i in ascii_out:
        for j in i:
            text_out += chr(10240 + int(j))
        text_out += "\n"

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    left = Inches(0.3)
    top = Inches(0.3)
    width = height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text_out
    p.font.size = Pt(FONT_SIZE)
    p.font.name = "Apple Braille"
    if frame % 100 == 0 and frame != 0:
        secLeft = round((time.time() - startTime) / 100 * (count - frame), 2)
        startTime = time.time()
        print("Written", frame, "frames, Remaining:", secLeft, "secs")

prs.save("test.pptx")
